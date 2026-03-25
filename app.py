import streamlit as st
import pandas as pd
from pptx import Presentation
import os, shutil, copy, re, time
from datetime import datetime
import numpy as np

import PIL.Image
import PIL.ImageOps
from PIL import ImageDraw, ImageFont, ExifTags

# Pillow 호환성 (최신 버전 대응)
if not hasattr(PIL.Image, 'ANTIALIAS'):
    PIL.Image.ANTIALIAS = PIL.Image.LANCZOS

from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip, clips_array
from moviepy.audio.fx.all import audio_loop, audio_fadeout

# ==========================================
# [도우미 함수] 1단계: PPT 복제 및 텍스트 교체
# ==========================================
def duplicate_slide(prs, source_slide):
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in new_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)
    for shape in source_slide.shapes:
        new_shape = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(copy.deepcopy(shape.element))
    return new_slide

def replace_text_in_shape(shape, keywords, row):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)
        needs_replacement = False
        for kw in keywords:
            if f"{{{{{kw}}}}}" in full_text:
                needs_replacement = True
                break
        if needs_replacement:
            for kw in keywords:
                target_text = f"{{{{{kw}}}}}"
                val = str(row[kw]) if pd.notna(row[kw]) else ""
                full_text = full_text.replace(target_text, val)
            if paragraph.runs:
                paragraph.runs[0].text = full_text
                for i in range(len(paragraph.runs) - 1, 0, -1):
                    p = paragraph._p
                    p.remove(paragraph.runs[i]._r)

# ==========================================
# [도우미 함수] 2단계: 사진 삽입 및 삭제
# ==========================================
def insert_photos_and_clean_ppt_com(ppt_path, photo_folder):
    import win32com.client, pythoncom
    pythoncom.CoInitialize() 
    abs_ppt_path = os.path.abspath(ppt_path)
    abs_photo_folder = os.path.abspath(photo_folder)
    pattern_regex = re.compile(r'\d{2}-\d{1}-\d{3}')
    
    inserted_count = 0
    deleted_management_numbers = []
    
    powerpoint = None
    presentation = None
    
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1 
        presentation = powerpoint.Presentations.Open(abs_ppt_path)
        
        for i in range(presentation.Slides.Count, 0, -1):
            slide = presentation.Slides.Item(i)
            management_number = None
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    text = shape.TextFrame.TextRange.Text
                    match = pattern_regex.search(text)
                    if match:
                        management_number = match.group()
                        break 
            
            if management_number:
                image_name = f"{management_number}.jpg"
                image_full_path = os.path.join(abs_photo_folder, image_name)
                if os.path.exists(image_full_path):
                    slide.Shapes.AddPicture(
                        FileName=image_full_path, LinkToFile=False, SaveWithDocument=True, 
                        Left=0, Top=175, Width=1290, Height=1150
                    )
                    inserted_count += 1
                else:
                    deleted_management_numbers.append(management_number)
                    slide.Delete() 
        
        today_str = datetime.now().strftime("%Y%m%d_%H%M%S")
        new_ppt_name = f"작업완료용_PPT_사진포함_{today_str}.pptx"
        new_ppt_path = os.path.abspath(os.path.join(os.getcwd(), new_ppt_name))
        
        presentation.SaveAs(new_ppt_path)
        presentation.Close()
        powerpoint.Quit()
        deleted_management_numbers.reverse()
        
        return True, new_ppt_path, inserted_count, deleted_management_numbers
        
    except Exception as e:
        if presentation:
            try: presentation.Close()
            except: pass
        if powerpoint:
            try: powerpoint.Quit()
            except: pass
        return False, str(e), 0, []
    finally:
        pythoncom.CoUninitialize()

# ==========================================
# [도우미 함수] 3단계: 이미지 추출 및 영상 제작
# ==========================================
def extract_images_from_ppt_regex(ppt_path, output_folder):
    import win32com.client, pythoncom
    pythoncom.CoInitialize() # 스레드 초기화
    try:
        abs_ppt, abs_out = os.path.abspath(ppt_path), os.path.abspath(output_folder)
        if os.path.exists(abs_out): shutil.rmtree(abs_out)
        os.makedirs(abs_out)
        pattern_regex = re.compile(r'\d{2}-\d{1}-\d{3}')
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        presentation = powerpoint.Presentations.Open(abs_ppt, WithWindow=False)
        extracted_count = 0
        for i, slide in enumerate(presentation.Slides):
            m_num = None
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    match = pattern_regex.search(shape.TextFrame.TextRange.Text)
                    if match:
                        m_num = match.group()
                        break
            img_name = f"{m_num}.jpg" if m_num else f"패턴없음_슬라이드{i+1:02d}.jpg"
            slide.Export(os.path.join(abs_out, img_name), "JPG")
            extracted_count += 1
            time.sleep(0.05) # RPC 에러 방지용 대기
        time.sleep(0.5)
        presentation.Close()
        powerpoint.Quit()
        return True, extracted_count
    except Exception as e:
        return False, str(e)
    finally:
        pythoncom.CoUninitialize()

# (NEW) 이미지 로드 함수 업데이트 (Pillow 객체 반환)
def load_fixed_image_pill(img_path, target_w, target_h):
    with PIL.Image.open(img_path) as img:
        img = PIL.ImageOps.exif_transpose(img).convert('RGB')
        img = img.resize((target_w, target_h), PIL.Image.LANCZOS)
        return img

# (NEW) 자막 오버레이 함수 분리
def apply_subtitle_draw(pill_img, custom_text):
    draw = ImageDraw.Draw(pill_img)
    w, h = pill_img.size
    try: font = ImageFont.truetype("malgun.ttf", 30)
    except: font = ImageFont.load_default()
    custom_text = (custom_text or "").strip()
    if custom_text:
        try: text_w = draw.textbbox((0, 0), custom_text, font=font)[2]
        except: text_w = draw.textsize(custom_text, font=font)[0]
        draw.text((w - text_w - 40, 40), custom_text, font=font, fill="black")
    return pill_img

# (업데이트) 3단계: 기존 렌더링 로직 함수화 (Standard 16:9)
def create_smart_video_standard(img_folder, audio_path, out_file, custom_text, img_dur, fps, trans_time):
    img_files = sorted([os.path.join(img_folder, f) for f in os.listdir(img_folder) if f.lower().endswith(('.jpg', '.jpeg', '.png'))])
    if not img_files: raise ValueError("이미지가 없습니다!")
    
    raw_clips, portrait_queue, first_portrait = [], [], None
    w_full, h_full = 1920, 1080 # 가로 FHD
    w_half, h_half = 960, 1080 # 가로 FHD 반

    for img_path in img_files:
        with PIL.Image.open(img_path) as img:
            w, h = PIL.ImageOps.exif_transpose(img).size
        if w > h:
            # 가로 사진: 꽉 채우기
            pill_img = load_fixed_image_pill(img_path, w_full, h_full)
            pill_img = apply_subtitle_draw(pill_img, custom_text)
            raw_clips.append(ImageClip(np.array(pill_img)).set_duration(img_dur))
        else:
            # 세로 사진: 큐에 저장 후 2개 병합
            if not first_portrait: first_portrait = img_path
            portrait_queue.append(img_path)
            if len(portrait_queue) == 2:
                pill1 = load_fixed_image_pill(portrait_queue[0], w_half, h_half)
                pill2 = load_fixed_image_pill(portrait_queue[1], w_half, h_half)
                combined_img = PIL.Image.new('RGB', (w_full, h_full))
                combined_img.paste(pill1, (0, 0))
                combined_img.paste(pill2, (w_half, 0))
                combined_img = apply_subtitle_draw(combined_img, custom_text)
                raw_clips.append(ImageClip(np.array(combined_img)).set_duration(img_dur))
                portrait_queue.clear()
    
    # 세로 사진 홀수 개 처리
    if len(portrait_queue) == 1:
        leftover = portrait_queue[0]
        partner = first_portrait if (first_portrait and first_portrait != leftover) else leftover
        pill1 = load_fixed_image_pill(leftover, w_half, h_half)
        pill2 = load_fixed_image_pill(partner, w_half, h_half)
        combined_img = PIL.Image.new('RGB', (w_full, h_full))
        combined_img.paste(pill1, (0, 0))
        combined_img.paste(pill2, (w_half, 0))
        combined_img = apply_subtitle_draw(combined_img, custom_text)
        raw_clips.append(ImageClip(np.array(combined_img)).set_duration(img_dur))

    # 영상 병합 로직 호출
    return merge_clips_and_audio(raw_clips, audio_path, out_file, fps, trans_time)

# (NEW) 3단계: 쇼츠 렌더링 로직 (Shorts 9:16)
def create_smart_video_shorts(img_folder, audio_path, out_file, custom_text, img_dur, fps, trans_time):
    img_files = sorted([os.path.join(img_folder, f) for f in os.listdir(img_folder) if f.lower().endswith(('.jpg', '.jpeg', '.png'))])
    if not img_files: raise ValueError("이미지가 없습니다!")
    
    raw_clips = []
    w_shorts, h_shorts = 1080, 1920 # 세로 FHD
    pad_top, pad_bottom = 250, 250 # 💡 쇼츠 정보 보호를 위한 위아래 여백 (Padding)
    content_h = h_shorts - pad_top - pad_bottom # 실제 콘텐츠 높이

    for img_path in img_files:
        with PIL.Image.open(img_path) as img:
            img_w, img_h = PIL.ImageOps.exif_transpose(img).size
            aspect_ratio = img_w / img_h

        if aspect_ratio >= (w_shorts / content_h):
            # 콘텐츠 영역에 비해 가로가 넓은 사진: 가로 맞춤
            target_w = w_shorts
            target_h = int(w_shorts / aspect_ratio)
        else:
            # 콘텐츠 영역에 비해 세로가 긴 사진: 세로 맞춤
            target_w = int(content_h * aspect_ratio)
            target_h = content_h

        # 콘텐츠 영역 내에서 크기 조정
        img_pill = load_fixed_image_pill(img_path, target_w, target_h)
        
        # 쇼츠용 배경 생성 (전체 화면, 검정색) 및 콘텐츠 배치
        combined_img = PIL.Image.new('RGB', (w_shorts, h_shorts), (0, 0, 0))
        pill_pos_x = (w_shorts - target_w) // 2
        pill_pos_y = pad_top + (content_h - target_h) // 2
        combined_img.paste(img_pill, (pill_pos_x, pill_pos_y))
        
        # 자막 오버레이
        combined_img = apply_subtitle_draw(combined_img, custom_text)
        raw_clips.append(ImageClip(np.array(combined_img)).set_duration(img_dur))

    # 영상 병합 로직 호출
    return merge_clips_and_audio(raw_clips, audio_path, out_file, fps, trans_time)

# (NEW) 클립 및 오디오 병합 공통 로직
def merge_clips_and_audio(raw_clips, audio_path, out_file, fps, trans_time):
    final_clips = [c.crossfadein(trans_time) if i > 0 else c for i, c in enumerate(raw_clips)]
    final_video = concatenate_videoclips(final_clips, padding=-trans_time, method="compose")

    if audio_path and os.path.exists(audio_path):
        audio = AudioFileClip(audio_path)
        audio = audio.subclip(0, final_video.duration) if audio.duration > final_video.duration else audio_loop(audio, duration=final_video.duration)
        final_video = final_video.set_audio(audio.fx(audio_fadeout, 2))

    final_video.write_videofile(out_file, fps=fps, codec="libx264", audio_codec="aac", bitrate="1000k", audio_bitrate="64k", preset="slow")
    final_video.close()
    if audio_path and os.path.exists(audio_path): audio.close()
    return out_file
# ==========================================
# [새로운 기능] 단계별 진행 상황 표시 바 (Progress Bar)
# ==========================================
def draw_progress_bar(current_step):
    c1 = "#A1E6AC" if current_step > 1 else ("#FFFFFF" if current_step == 1 else "#E0E0E0")
    c2 = "#A1E6AC" if current_step > 2 else ("#FFFFFF" if current_step == 2 else "#E0E0E0")
    c3 = "#FFFFFF" if current_step == 3 else "#E0E0E0"

    t1 = "white" if current_step > 1 else ("black" if current_step == 1 else "#888888")
    t2 = "white" if current_step > 2 else ("black" if current_step == 2 else "#888888")
    t3 = "black" if current_step >= 3 else "#888888"

    i1 = "✅" if current_step > 1 else "📝"
    i2 = "✅" if current_step > 2 else "📸"
    i3 = "🎬"

    border_color = "#E0E0E0"

    html = f"""
    <div style='display: flex; justify-content: space-between; align-items: center; margin-bottom: 30px; margin-top: 10px; font-family: "Malgun Gothic", sans-serif;'>
        <div style='flex: 1; text-align: center; background-color: {c1}; color: {t1}; padding: 20px; border-radius: 15px; font-size: 20px; font-weight: bold; border: 2px solid {border_color}; transition: 0.3s;'>
            <span style='font-size: 24px;'>{i1}</span><br>Step 1. 텍스트 PPT 생성
        </div>
        <div style='font-size: 40px; margin: 0 20px; color: {border_color};'>➔</div>
        <div style='flex: 1; text-align: center; background-color: {c2}; color: {t2}; padding: 20px; border-radius: 15px; font-size: 20px; font-weight: bold; border: 2px solid {border_color}; transition: 0.3s;'>
            <span style='font-size: 24px;'>{i2}</span><br>Step 2. 사진 자동 삽입
        </div>
        <div style='font-size: 40px; margin: 0 20px; color: {border_color};'>➔</div>
        <div style='flex: 1; text-align: center; background-color: {c3}; color: {t3}; padding: 20px; border-radius: 15px; font-size: 20px; font-weight: bold; border: 2px solid {border_color}; transition: 0.3s;'>
            <span style='font-size: 24px;'>{i3}</span><br>Step 3. 영상 렌더링
        </div>
    </div>
    """
    st.markdown(html, unsafe_allow_html=True)
# ==========================================
# 🌟 Streamlit UI 설정 및 디자인 적용
# ==========================================

# ==========================================
# 🌟 Streamlit UI 설정 및 전역 디자인 적용
# ==========================================
st.set_page_config(page_title="동영상 제작 자동화", layout="wide")

# [Rounded Rectangular Title Bar]
title_bar_html = """
<div style='background-color: #FFFFFF; padding: 20px; border-radius: 20px; margin-bottom: 20px; text-align: center; border: 2px solid #E0E0E0; font-family: "Malgun Gothic", sans-serif;'>
    <h1 style='color: #000000; margin: 0;'>🐶 동영상 제작 통합 자동화</h1>
</div>
"""
st.markdown(title_bar_html, unsafe_allow_html=True)

# (NEW) 파스텔 톤 CSS 주입 (버튼, 텍스트, 위젯 느낌 개선 및 탭 글자 크기 확대)
custom_css = """
<style>
    /* 🌟 상단 탭(Tab) 메뉴 글자 크기 및 굵기 변경 🌟 */
    button[data-baseweb="tab"] > div[data-testid="stMarkdownContainer"] > p {
        font-size: 22px !important; /* 👈 여기서 글자 크기를 조절하세요 (기본은 약 16px) */
        font-weight: 700 !important; /* 글자 굵기 (숫자가 클수록 굵어짐) */
    }
    
    /* 탭 메뉴 높이 여백 조정 (글자가 커졌으므로 공간을 살짝 넓혀줍니다) */
    button[data-baseweb="tab"] {
        padding-top: 15px !important;
        padding-bottom: 15px !important;
    }

    /* 전체 배경을 아주 연한 아이보리 계열로 살짝 조정 (선택사항) */
    .stApp {
        background-color: #FAFAFA;
    }
    
    /* 기본 버튼 디자인 (연한 노란색 + 오렌지색 호버) */
    div.stButton > button:first-child {
        background-color: #FFF9C4;
        color: #333333;
        border: 1px solid #FFB74D;
        border-radius: 8px;
        font-weight: 600;
        transition: all 0.3s ease;
    }
    div.stButton > button:first-child:hover {
        background-color: #FFB74D;
        color: white;
        border-color: #FFB74D;
        box-shadow: 0 4px 8px rgba(255, 183, 77, 0.4);
    }
    
    /* Primary 버튼 디자인 (강조용 - 오렌지 톤) */
    div.stButton > button[data-baseweb="button"][kind="primary"] {
        background-color: #FF9800;
        color: white;
        border: none;
    }
    div.stButton > button[data-baseweb="button"][kind="primary"]:hover {
        background-color: #F57C00;
    }
    
    /* 헤더 스타일링 */
    h1, h2, h3 {
        color: #424242;
    }
</style>
"""
st.markdown(custom_css, unsafe_allow_html=True)
# ==========================================
# 🌟 작업 디렉토리 및 세션 상태 초기화
# ==========================================
DIR_EXTRACTED, DIR_SELECTED = os.path.abspath("temp_extracted"), os.path.abspath("temp_selected")
for d in [DIR_EXTRACTED, DIR_SELECTED]:
    if not os.path.exists(d): os.makedirs(d)

for key in ['step1_ppt_path', 'step2_ppt_path']:
    if key not in st.session_state: st.session_state[key] = None
for key in ['images_extracted', 'images_selected']:
    if key not in st.session_state: st.session_state[key] = False


# ==========================================
# 🌟 탭(Tab) 기반 UI 구성
# ==========================================
tabs = st.tabs(["📝 1단계: 텍스트 PPT 생성", "📸 2단계: 사진 삽입 & 정리", "🎬 3단계: 영상 렌더링"])

# --- [1단계 탭] ---
with tabs[0]:
    draw_progress_bar(1) 
    st.header("Step 1: 엑셀 데이터로 텍스트 뼈대 PPT 생성")
    col1, col2 = st.columns(2)
    with col1:
        excel_file = st.file_uploader("1. 데이터 업로드 (Excel)", type=['xlsx'])
        pptx_template = st.file_uploader("2. 템플릿 업로드 (PPTX)", type=['pptx'])
    with col2:
        if st.button("🏗️ Step 1: 텍스트 PPT 생성 시작", use_container_width=True):
            if excel_file and pptx_template:
                with st.spinner("텍스트 데이터를 PPT에 입력 중입니다..."):
                    try:
                        df = pd.read_excel(excel_file)
                        df.columns = df.columns.str.strip()
                        keywords = ["관리번호", "견종", "성별", "나이", "몸무게"]
                        prs = Presentation(pptx_template)
                        source_slide = prs.slides[0]
                        slides_to_edit = [source_slide]
                        for _ in range(1, len(df)):
                            slides_to_edit.append(duplicate_slide(prs, source_slide))
                        for i, row in df.iterrows():
                            current_slide = slides_to_edit[i]
                            for shape in current_slide.shapes:
                                replace_text_in_shape(shape, keywords, row)
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        output_ppt_path = f"사진인화작업용_PPT_no이미지_{timestamp}.pptx"
                        prs.save(output_ppt_path)
                        st.session_state['step1_ppt_path'] = output_ppt_path
                        st.success(f"✔️ Step 1 완료! ({output_ppt_path})")
                    except Exception as e:
                        st.error(f"오류 발생: {e}")
            else:
                st.warning("엑셀 파일과 템플릿을 업로드해주세요.")

# --- [2단계 탭] ---
with tabs[1]:
    draw_progress_bar(2) 
    st.header("Step 2: 사진 일괄 삽입 및 미존재 슬라이드 정리")
    ppt_target = st.session_state.get('step1_ppt_path')
    if not ppt_target:
        st.warning("⚠️ 1단계에서 PPT를 먼저 생성해야 합니다.")
    else:
        st.info(f"📄 현재 작업 대상 파일: **{ppt_target}**")
        photo_dir_input = st.text_input("📁 사진이 들어있는 폴더의 절대 경로를 입력하세요. 🚨해당 폴더가 열려 있으면 에러 발생")
        if st.button("📸 Step 2: 사진 삽입 및 슬라이드 정리 시작", use_container_width=True):
            if not photo_dir_input or not os.path.exists(photo_dir_input.strip('"').strip("'")):
                st.error("정확한 사진 폴더 경로를 입력해주세요.")
            else:
                clean_photo_dir = photo_dir_input.strip('"').strip("'")
                with st.spinner("사진 삽입 및 슬라이드 정리 중... (파워포인트가 실행됩니다)"):
                    success, result_path, in_count, del_list = insert_photos_and_clean_ppt_com(ppt_target, clean_photo_dir)
                    if success:
                        st.session_state['step2_ppt_path'] = result_path
                        st.success(f"🎉 총 {in_count}장의 사진 삽입 완료!")
                        if del_list:
                            with st.expander(f"🚨 사진이 없어 삭제된 슬라이드 ({len(del_list)}개)", expanded=True):
                                st.error(f"사진 누락으로 삭제된 관리번호:\n\n**{', '.join(del_list)}**")
                        st.success(f"💾 최종 PPT 저장 완료: {result_path}")
                    else:
                        st.error(f"오류 발생: {result_path}")

# --- [3단계 탭] ---
with tabs[2]:
    draw_progress_bar(3) 
    st.header("Step 3: 스마트 이미지 추출 및 영상 렌더링")
    
    # 💡 [핵심 기능] 테스트 모드 토글 스위치
    skip_mode = st.toggle("🚀 [테스트 모드] 1, 2단계를 건너뛰고 사진 폴더로 바로 영상 렌더링 테스트하기")
    
    if skip_mode:
        st.info("💡 PPT 작업을 생략하고 지정한 폴더의 사진을 바로 영상으로 만듭니다.")
        test_photo_dir = st.text_input("📁 테스트할 원본 사진들이 들어있는 폴더의 절대 경로를 입력하세요")
        
        col_t1, col_t2 = st.columns(2)
        with col_t1:
            if st.button("🔍 이 사진들로 '3-2. 엑셀 선별'부터 테스트", use_container_width=True):
                clean_dir = test_photo_dir.strip('"').strip("'")
                if not os.path.exists(clean_dir):
                    st.error("폴더 경로를 찾을 수 없습니다.")
                else:
                    with st.spinner("사진 불러오는 중..."):
                        if os.path.exists(DIR_EXTRACTED): shutil.rmtree(DIR_EXTRACTED)
                        os.makedirs(DIR_EXTRACTED)
                        for f in os.listdir(clean_dir):
                            if f.lower().endswith(('.jpg', '.jpeg', '.png')):
                                shutil.copy2(os.path.join(clean_dir, f), os.path.join(DIR_EXTRACTED, f))
                        st.session_state['images_extracted'] = True
                        st.success("✔️ 사진 로드 완료! 아래 '3-2'에서 엑셀을 업로드해 선별하세요.")
        with col_t2:
            if st.button("🎬 엑셀 선별 없이 '3-3. 영상 렌더링' 직행", use_container_width=True):
                clean_dir = test_photo_dir.strip('"').strip("'")
                if not os.path.exists(clean_dir):
                    st.error("폴더 경로를 찾을 수 없습니다.")
                else:
                    with st.spinner("사진 불러오는 중..."):
                        if os.path.exists(DIR_SELECTED): shutil.rmtree(DIR_SELECTED)
                        os.makedirs(DIR_SELECTED)
                        for f in os.listdir(clean_dir):
                            if f.lower().endswith(('.jpg', '.jpeg', '.png')):
                                shutil.copy2(os.path.join(clean_dir, f), os.path.join(DIR_SELECTED, f))
                        st.session_state['images_extracted'] = True
                        st.session_state['images_selected'] = True
                        st.success("✔️ 사진 로드 완료! 아래 '3-3'에서 바로 영상을 렌더링하세요.")
        st.divider()

    final_ppt = st.session_state.get('step2_ppt_path')
    
    if not skip_mode and not final_ppt:
        st.warning("⚠️ 2단계에서 최종 PPT를 먼저 생성해주세요. (또는 위 '테스트 모드' 토글을 켜세요)")
    
    # [3-1] 이미지 추출 (테스트 모드가 아닐 때만 노출)
    if not skip_mode and final_ppt:
        st.subheader("3-1. 슬라이드 이미지 추출")
        if st.button("🖼️ 1. PPT를 이미지로 추출 (클릭)"):
            with st.spinner("관리번호 패턴을 분석하여 슬라이드를 이미지로 변환 중입니다..."):
                success, count_or_err = extract_images_from_ppt_regex(final_ppt, DIR_EXTRACTED)
                if success:
                    st.session_state['images_extracted'] = True
                    st.success(f"✔️ 총 {count_or_err}장의 이미지 추출 완료! (임시 폴더 저장됨)")
                else:
                    st.error(f"추출 실패: {count_or_err}")

    # [3-2] 엑셀 기반 선별
    if st.session_state.get('images_extracted'):
        if not skip_mode: st.divider()
        st.subheader("3-2. 영상 제작용 데이터 선별")
        st.write("💡 '제작대상' 열에 'O'가 표시된 관리번호만 영상 재료로 선별합니다.")
        filter_excel = st.file_uploader("선별용 엑셀 업로드", type=['xlsx'], key='filter_xls')
        
        if filter_excel:
            if st.button("🔍 2. 선별 데이터 바탕으로 이미지 복사"):
                with st.spinner("제작 대상 이미지를 추려내는 중..."):
                    try:
                        if os.path.exists(DIR_SELECTED): shutil.rmtree(DIR_SELECTED)
                        os.makedirs(DIR_SELECTED)
                        
                        df_filter = pd.read_excel(filter_excel)
                        targets = df_filter[df_filter['제작대상'].astype(str).str.upper() == 'O']['관리번호'].tolist()
                        
                        success_cnt, fail_cnt = 0, 0
                        for base_name in targets:
                            source_file = os.path.join(DIR_EXTRACTED, f"{base_name}.jpg")
                            if os.path.exists(source_file):
                                shutil.copy2(source_file, os.path.join(DIR_SELECTED, f"{base_name}.jpg"))
                                success_cnt += 1
                            else:
                                fail_cnt += 1
                                
                        st.session_state['images_selected'] = True
                        st.success(f"✔️ 선별 완료: 복사 성공 {success_cnt}건, 실패(이미지 없음) {fail_cnt}건")
                    except Exception as e:
                        st.error(f"선별 중 오류 발생: {e}")

    # [3-3] 동영상 제작 설정 및 렌더링
    if st.session_state.get('images_selected'):
        st.divider()
        st.subheader("3-3. 스마트 동영상 렌더링")
        colA, colB, colC = st.columns([1, 1.5, 1])
        with colA:
            video_ratio = st.selectbox(
                "🎥 출력 동영상 비율 선택",
                ("Standard (16:9 Landscape)", "Shorts (9:16 Vertical)")
            )
        with colB:
            custom_subtitle = st.text_input("우측 상단 자막 입력", "입양 문의: 042-270-7239")
            bgm_file = st.file_uploader("배경음악 (BGM) 업로드", type=['mp3', 'wav'])
        with colC:
            img_duration = st.slider("사진 1장당 재생 시간(초)", 2, 10, 4)
            fps_val = st.slider("FPS (초당 프레임)", 10, 60, 15)
            transition_val = st.slider("전환 효과 시간(초)", 0, 3, 1)

        if st.button("🎬 3. 최적화 동영상 제작 시작", type="primary", use_container_width=True):
            with st.spinner("초고화질 최적화 인코딩을 시작합니다. 잠시만 기다려주세요..."):
                try:
                    bgm_path = None
                    if bgm_file:
                        bgm_path = "temp_bgm_file.mp3"
                        with open(bgm_path, "wb") as f:
                            f.write(bgm_file.read())

                    if video_ratio == "Shorts (9:16 Vertical)":
                        video_name_prefix = "최종완성_쇼츠영상"
                        video_p_func = create_smart_video_shorts
                    else:
                        video_name_prefix = "최종완성_가로영상"
                        video_p_func = create_smart_video_standard

                    output_video_name = f"{video_name_prefix}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.mp4"
                    
                    output_final = video_p_func(
                        img_folder=DIR_SELECTED, 
                        audio_path=bgm_path, 
                        out_file=output_video_name,
                        custom_text=custom_subtitle,
                        img_dur=img_duration,          
                        fps=fps_val,
                        trans_time=transition_val
                    )
                    
                    st.success(f"🎉 성공! {video_ratio} 형식이 완벽하게 렌더링되었습니다!")
                    st.video(output_final)
                    
                    if bgm_path and os.path.exists(bgm_path):
                        os.remove(bgm_path)
                        
                except Exception as e:
                    st.error(f"동영상 렌더링 중 오류 발생: {e}")